import os
import re
import sys
import zipfile
import unicodedata
import xml.etree.ElementTree as ET
import fitz # PyMuPDF
import docx
import pptx
import win32com.client
from docx.text.paragraph import Paragraph
from docx.table import Table

# Configure output encoding to utf-8
sys.stdout.reconfigure(encoding='utf-8')

# Input / Output Directories
SRC_DIR = r"D:\Users\Admin\Downloads\cảm biến và đo lường cho robot"
OUT_MD_DIR = r"C:\Users\Admin\.gemini\antigravity\scratch\2526II_RBE3042_1\docs\lectures"
OUT_FIG_DIR = r"C:\Users\Admin\.gemini\antigravity\scratch\2526II_RBE3042_1\figures\lectures"

# Ensure directories exist
os.makedirs(OUT_MD_DIR, exist_ok=True)
os.makedirs(OUT_FIG_DIR, exist_ok=True)

# Target files list
TARGET_FILES = [
    "EP10 Pressure.pdf",
    "Chappter 5.2 ADC-Sampling and Filter.pptx",
    "EP5 Circuits for Sensors.pdf",
    "EP5b Interface Electronic Circuits.pdf",
    "Chương 4.docx",
    "EP2 Measurement Systems.pdf",
    "EP3 Errors during the measurement process.pdf",
    "EP4 Variable conversion elements.pdf",
    "de-cuong-on-tap-ky-thuat-do-luong-va-cam-bien-k62-uet-450.pdf",
    "AD620.pdf"
]

def normalize_filename(filename):
    """Normalize string to remove accents and lowercase for flexible matching."""
    return unicodedata.normalize('NFKD', filename).encode('ASCII', 'ignore').decode('ASCII').lower()

def find_actual_path(target_filename):
    """Finds the actual file path in SRC_DIR using flexible matching."""
    if not os.path.exists(SRC_DIR):
        raise FileNotFoundError(f"Source directory not found: {SRC_DIR}")
    
    # Try exact match first
    exact_path = os.path.join(SRC_DIR, target_filename)
    if os.path.exists(exact_path):
        return exact_path
        
    # Try normalized match
    norm_target = normalize_filename(target_filename)
    for f in os.listdir(SRC_DIR):
        if normalize_filename(f) == norm_target:
            return os.path.join(SRC_DIR, f)
            
    # Try substring match
    for f in os.listdir(SRC_DIR):
        norm_f = normalize_filename(f)
        if norm_target in norm_f or norm_f in norm_target:
            return os.path.join(SRC_DIR, f)
            
    raise FileNotFoundError(f"Could not find matching file for '{target_filename}' in {SRC_DIR}")


# ==========================================
# 1. PPTX Processing Function
# ==========================================
def process_pptx(pptx_path):
    pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]
    print(f"\n--- Processing PPTX: {pptx_name} ---")
    
    # parse slide relationship XMLs to map slide numbers to media files
    slide_media = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith('ppt/slides/_rels/slide') and name.endswith('.xml.rels'):
                basename = os.path.basename(name)
                slide_num_str = basename.replace('slide', '').replace('.xml.rels', '')
                try:
                    slide_num = int(slide_num_str)
                except ValueError:
                    continue
                
                xml_content = z.read(name)
                root = ET.fromstring(xml_content)
                
                media_files = []
                # ponytail: standard namespace matching to find Relationship
                for rel in root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                    target = rel.attrib.get('Target', '')
                    rel_id = rel.attrib.get('Id', '')
                    if 'media/' in target:
                        clean_target = target.replace('../', 'ppt/')
                        if clean_target in z.namelist():
                            media_files.append((clean_target, rel_id))
                if media_files:
                    slide_media[slide_num] = media_files

    # Extract text from slide shapes using python-pptx
    prs = pptx.Presentation(pptx_path)
    slide_texts = {}
    for idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        slide_texts[idx] = "\n".join(text_parts)

    # Use win32com to export slide as high-quality PNG
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    # Open presentation (FileName, ReadOnly, Untitled, WithWindow)
    pres = ppt_app.Presentations.Open(os.path.abspath(pptx_path), True, False, False)
    
    width_pts = pres.PageSetup.SlideWidth
    height_pts = pres.PageSetup.SlideHeight
    scale_w = int(width_pts * (150.0 / 72.0))
    scale_h = int(height_pts * (150.0 / 72.0))
    
    slide_images = {}
    for idx, slide in enumerate(pres.Slides, start=1):
        img_filename = f"{pptx_name}_slide_{idx}.png"
        img_path = os.path.join(OUT_FIG_DIR, img_filename)
        slide.Export(os.path.abspath(img_path), "PNG", scale_w, scale_h)
        slide_images[idx] = img_filename
        
    pres.Close()
    ppt_app.Quit()

    # Extract slide media files from zip to output folder
    slide_extracted_media = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for slide_num, media_list in slide_media.items():
            extracted = []
            for zip_path, rel_id in media_list:
                ext = os.path.splitext(zip_path)[1]
                out_filename = f"{pptx_name}_slide_{slide_num}_media_{rel_id}{ext}"
                out_path = os.path.join(OUT_FIG_DIR, out_filename)
                
                with open(out_path, 'wb') as f:
                    f.write(z.read(zip_path))
                extracted.append(out_filename)
            slide_extracted_media[slide_num] = extracted

    # Write Markdown output
    md_content = []
    md_content.append(f"# {pptx_name}\n\n")
    md_content.append(f"> Tài liệu chuyển đổi từ PPTX: `{os.path.basename(pptx_path)}`\n\n")
    md_content.append("---\n\n")
    
    for idx in range(1, len(prs.slides) + 1):
        md_content.append(f"## Slide {idx}\n\n")
        
        text = slide_texts.get(idx, "").strip()
        if text:
            for line in text.split('\n'):
                line = line.strip()
                if line:
                    md_content.append(f"- {line}\n")
            md_content.append("\n")
            
        slide_img = slide_images.get(idx)
        if slide_img:
            md_content.append(f"![Slide {idx}](../../figures/lectures/{slide_img})\n\n")
            
        extracted_media = slide_extracted_media.get(idx, [])
        for media_file in extracted_media:
            md_content.append(f"![Embedded Media](../../figures/lectures/{media_file})\n\n")
            
        md_content.append("---\n\n")
        
    md_path = os.path.join(OUT_MD_DIR, f"{pptx_name}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("".join(md_content))
        
    print(f"Done PPTX: {pptx_name}")


# ==========================================
# 2. DOCX Processing Function
# ==========================================
def process_docx(docx_path):
    docx_name = os.path.splitext(os.path.basename(docx_path))[0]
    print(f"\n--- Processing DOCX: {docx_name} ---")
    
    # Parse document relationships to map media files
    rels = {}
    with zipfile.ZipFile(docx_path, 'r') as z:
        if 'word/_rels/document.xml.rels' in z.namelist():
            xml_content = z.read('word/_rels/document.xml.rels')
            root = ET.fromstring(xml_content)
            for rel in root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                r_id = rel.attrib.get('Id')
                target = rel.attrib.get('Target')
                if r_id and target:
                    if target.startswith('media/'):
                        rels[r_id] = 'word/' + target
                    elif target.startswith('/word/media/'):
                        rels[r_id] = target.lstrip('/')
                    else:
                        rels[r_id] = 'word/' + target
                        
    def get_or_extract_media(r_id):
        if r_id not in rels:
            return None
        zip_path = rels[r_id]
        ext = os.path.splitext(zip_path)[1]
        out_filename = f"{docx_name}_{r_id}{ext}"
        out_path = os.path.join(OUT_FIG_DIR, out_filename)
        
        with zipfile.ZipFile(docx_path, 'r') as z:
            if zip_path in z.namelist():
                with open(out_path, 'wb') as f:
                    f.write(z.read(zip_path))
                return out_filename
        return None

    # Math ML (OMath) to LaTeX conversion function
    M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    
    def omath_to_latex(element):
        tag = element.tag.split('}')[-1]
        
        # Text run
        if tag == 'r':
            t_el = element.find(f'.//{{{W_NS}}}t') or element.find(f'.//{{{M_NS}}}t')
            if t_el is not None and t_el.text:
                return t_el.text
            return ""
            
        # Fraction
        elif tag == 'f':
            num = element.find(f'./{{{M_NS}}}num')
            den = element.find(f'./{{{M_NS}}}den')
            num_str = "".join([omath_to_latex(child) for child in num]) if num is not None else ""
            den_str = "".join([omath_to_latex(child) for child in den]) if den is not None else ""
            return f"\\frac{{{num_str}}}{{{den_str}}}"
            
        # Subscript
        elif tag in ['ssub', 'sSub']:
            e = element.find(f'./{{{M_NS}}}e')
            sub = element.find(f'./{{{M_NS}}}sub')
            e_str = "".join([omath_to_latex(child) for child in e]) if e is not None else ""
            sub_str = "".join([omath_to_latex(child) for child in sub]) if sub is not None else ""
            return f"{e_str}_{{{sub_str}}}"
            
        # Superscript
        elif tag in ['ssup', 'sSup']:
            e = element.find(f'./{{{M_NS}}}e')
            sup = element.find(f'./{{{M_NS}}}sup')
            e_str = "".join([omath_to_latex(child) for child in e]) if e is not None else ""
            sup_str = "".join([omath_to_latex(child) for child in sup]) if sup is not None else ""
            return f"{e_str}^{{{sup_str}}}"
            
        # Sub-Superscript
        elif tag in ['ssubsup', 'sSubSup']:
            e = element.find(f'./{{{M_NS}}}e')
            sub = element.find(f'./{{{M_NS}}}sub')
            sup = element.find(f'./{{{M_NS}}}sup')
            e_str = "".join([omath_to_latex(child) for child in e]) if e is not None else ""
            sub_str = "".join([omath_to_latex(child) for child in sub]) if sub is not None else ""
            sup_str = "".join([omath_to_latex(child) for child in sup]) if sup is not None else ""
            return f"{e_str}_{{{sub_str}}}^{{{sup_str}}}"
            
        # Delimiter
        elif tag == 'd':
            e_elements = element.findall(f'./{{{M_NS}}}e')
            sep_char = ""
            sep_el = element.find(f'./{{{M_NS}}}dPr/{{{M_NS}}}sepChr')
            if sep_el is not None and sep_el.get(f'{{{M_NS}}}val'):
                sep_char = sep_el.get(f'{{{M_NS}}}val')
                
            e_strs = []
            for e in e_elements:
                e_strs.append("".join([omath_to_latex(child) for child in e]))
                
            delimiter_content = sep_char.join(e_strs)
            
            beg_char = "("
            end_char = ")"
            beg_el = element.find(f'./{{{M_NS}}}dPr/{{{M_NS}}}beg')
            end_el = element.find(f'./{{{M_NS}}}dPr/{{{M_NS}}}end')
            if beg_el is not None and beg_el.get(f'{{{M_NS}}}val'):
                beg_char = beg_el.get(f'{{{M_NS}}}val')
            if end_el is not None and end_el.get(f'{{{M_NS}}}val'):
                end_char = end_el.get(f'{{{M_NS}}}val')
                
            beg_latex = f"\\left{beg_char}" if beg_char in ['(', '[', '{', '|'] else beg_char
            end_latex = f"\\right{end_char}" if end_char in [')', ']', '}', '|'] else end_char
            
            return f"{beg_latex}{delimiter_content}{end_latex}"
            
        # Radical
        elif tag == 'rad':
            e = element.find(f'./{{{M_NS}}}e')
            deg = element.find(f'./{{{M_NS}}}deg')
            e_str = "".join([omath_to_latex(child) for child in e]) if e is not None else ""
            if deg is not None and len(deg) > 0:
                deg_str = "".join([omath_to_latex(child) for child in deg])
                return f"\\sqrt[{deg_str}]{{{e_str}}}"
            return f"\\sqrt{{{e_str}}}"
            
        # Default recurse
        else:
            parts = []
            for child in element:
                parts.append(omath_to_latex(child))
            return "".join(parts)

    def traverse_element(elem):
        """Recursively traverse paragraph XML child elements sequentially."""
        tag = elem.tag.split('}')[-1]
        
        # Check for drawing/pict (image)
        if tag in ['drawing', 'pict']:
            rel_ids = []
            for desc in elem.iter():
                for attr_name, attr_val in desc.attrib.items():
                    if attr_name.endswith('}embed') or attr_name.endswith('}id'):
                        rel_ids.append(attr_val)
            
            md_imgs = []
            for r_id in rel_ids:
                filename = get_or_extract_media(r_id)
                if filename:
                    md_imgs.append(f"\n\n![Image](../../figures/lectures/{filename})\n\n")
            if md_imgs:
                return "".join(md_imgs)
            return ""
            
        # Check for math elements
        elif tag == 'oMath':
            latex = omath_to_latex(elem)
            return f" ${latex}$ "
            
        elif tag == 'oMathPara':
            latex = omath_to_latex(elem)
            return f"\n\n$$\n{latex}\n$$\n\n"
            
        elif tag == 't':
            return elem.text if elem.text else ""
            
        elif tag == 'tab':
            return "\t"
            
        elif tag in ['br', 'cr']:
            return "\n"
            
        # Recurse
        parts = []
        for child in elem:
            parts.append(traverse_element(child))
        return "".join(parts)

    def process_p(p):
        prefix = ""
        style_name = p.style.name.lower() if p.style else ""
        if style_name.startswith('heading 1'):
            prefix = "# "
        elif style_name.startswith('heading 2'):
            prefix = "## "
        elif style_name.startswith('heading 3'):
            prefix = "### "
        elif style_name.startswith('heading'):
            prefix = "#### "
        elif 'bullet' in style_name or 'list' in style_name:
            prefix = "- "
            
        content = traverse_element(p._element)
        if prefix:
            return f"{prefix}{content.strip()}\n\n"
        return f"{content}\n"

    # Read paragraphs and tables sequentially
    doc = docx.Document(docx_path)
    md_content = []
    md_content.append(f"# {docx_name}\n\n")
    md_content.append(f"> Tài liệu chuyển đổi từ DOCX: `{os.path.basename(docx_path)}`\n\n")
    md_content.append("---\n\n")

    # Iterate body block items in order
    def iter_block_items(parent):
        if isinstance(parent, docx.document.Document):
            parent_elm = parent.element.body
        elif isinstance(parent, docx.table._Cell):
            parent_elm = parent._tc
        else:
            raise TypeError("Unknown parent type")
        for child in parent_elm.iterchildren():
            if child.tag.endswith('p'):
                yield Paragraph(child, parent)
            elif child.tag.endswith('tbl'):
                yield Table(child, parent)

    for item in iter_block_items(doc):
        if isinstance(item, Paragraph):
            md_content.append(process_p(item))
        elif isinstance(item, Table):
            md_content.append("\n")
            for r_idx, row in enumerate(item.rows):
                row_cells = []
                for cell in row.cells:
                    cell_text = "".join([process_p(p).strip() for p in cell.paragraphs])
                    cell_text = cell_text.replace('\n', ' ')
                    row_cells.append(cell_text)
                md_content.append("| " + " | ".join(row_cells) + " |\n")
                if r_idx == 0:
                    md_content.append("| " + " | ".join(['---'] * len(row_cells)) + " |\n")
            md_content.append("\n")

    # Save output markdown
    md_path = os.path.join(OUT_MD_DIR, f"{docx_name}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        raw_md = "".join(md_content)
        # ponytail: clean redundant empty lines in output markdown
        clean_md = re.sub(r'\n{3,}', '\n\n', raw_md)
        f.write(clean_md)
        
    print(f"Done DOCX: {docx_name}")


# ==========================================
# 3. PDF Processing Function
# ==========================================
def process_pdf(pdf_path):
    pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
    print(f"\n--- Processing PDF: {pdf_name} ---")
    
    doc = fitz.open(pdf_path)
    md_content = []
    md_content.append(f"# {pdf_name}\n\n")
    md_content.append(f"> Tài liệu chuyển đổi từ PDF: `{os.path.basename(pdf_path)}`\n\n")
    md_content.append("---\n\n")
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        page_idx = page_num + 1
        
        # 1. Extract text
        text = page.get_text()
        
        # 2. Render page to high-quality PNG (DPI >= 150)
        # 150 DPI = 150/72 scale factor
        zoom = 150.0 / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        page_img_filename = f"{pdf_name}_page_{page_idx}.png"
        page_img_path = os.path.join(OUT_FIG_DIR, page_img_filename)
        pix.save(page_img_path)
        
        # 3. Extract embedded images, filtering out the 5895-byte logo
        image_list = page.get_images(full=True)
        extracted_images = []
        
        for img_idx, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            # Check for redundant template logo (exactly 5895 bytes)
            if len(image_bytes) == 5895:
                continue
                
            img_filename = f"{pdf_name}_page_{page_idx}_img_{img_idx}.{image_ext}"
            img_path = os.path.join(OUT_FIG_DIR, img_filename)
            with open(img_path, "wb") as f:
                f.write(image_bytes)
            extracted_images.append(img_filename)
            
        # 4. Formulate markdown structure
        md_content.append(f"## Trang {page_idx}\n\n")
        
        if text.strip():
            for line in text.strip().split('\n'):
                line = line.strip()
                if line:
                    md_content.append(f"- {line}\n")
            md_content.append("\n")
            
        md_content.append(f"![Trang {page_idx}](../../figures/lectures/{page_img_filename})\n\n")
        
        for img_file in extracted_images:
            md_content.append(f"![Hình ảnh trang {page_idx}](../../figures/lectures/{img_file})\n\n")
            
        md_content.append("---\n\n")
        
    md_path = os.path.join(OUT_MD_DIR, f"{pdf_name}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("".join(md_content))
        
    print(f"Done PDF: {pdf_name}")


# ==========================================
# Main Execution Driver
# ==========================================
def main():
    print("Starting Lecture Digitization script (R1)...")
    
    for filename in TARGET_FILES:
        try:
            actual_path = find_actual_path(filename)
            print(f"\nResolved '{filename}' -> '{actual_path}'")
            
            ext = os.path.splitext(actual_path)[1].lower()
            if ext == '.pdf':
                process_pdf(actual_path)
            elif ext == '.pptx':
                process_pptx(actual_path)
            elif ext == '.docx':
                process_docx(actual_path)
            else:
                print(f"Unsupported file type: {ext} for {filename}")
                
        except Exception as e:
            print(f"ERROR processing '{filename}': {e}", file=sys.stderr)
            import traceback
            traceback.print_exc()

    # Clean up any existing 5895-byte logo files in OUT_FIG_DIR
    print("\nCleaning up 5895-byte logo files from figures folder...")
    for f in os.listdir(OUT_FIG_DIR):
        path = os.path.join(OUT_FIG_DIR, f)
        if os.path.isfile(path) and os.path.getsize(path) == 5895:
            try:
                os.remove(path)
                print(f"Removed redundant logo file: {f}")
            except Exception as e:
                print(f"Error removing {f}: {e}")

if __name__ == "__main__":
    main()
